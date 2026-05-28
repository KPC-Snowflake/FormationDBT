"""
Genere le PPTX de la formation dbt a partir du template KPC.

Identite visuelle KPC :
- Font: Montserrat
- Couleurs: dk1=#160D3D, accent1=#320CC9, accent2=#FF660D
- Slide 16:9 (13.3" x 7.5")
- Layout "Vide" : placeholder idx=15 en haut (y=0, h=760021) = barre de titre
"""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

TEMPLATE = os.path.join(os.path.dirname(__file__), "Template.pptx")
OUTPUT = os.path.join(os.path.dirname(__file__), "Formation_dbt_Snowflake_KPC.pptx")

DK1 = RGBColor(0x16, 0x0D, 0x3D)
ACCENT1 = RGBColor(0x32, 0x0C, 0xC9)
ACCENT2 = RGBColor(0xFF, 0x66, 0x0D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)
FONT = "Montserrat"

TITLE_Y_BOTTOM = 850000
MARGIN_LEFT = 700000
CONTENT_TOP = TITLE_Y_BOTTOM + 50000


def add_textbox(slide, left, top, width, height, text, size=14,
                bold=False, color=DK1, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(4)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = FONT
    return txBox, tf


def add_para(tf, text, size=14, bold=False, color=DK1, space_before=0, level=0):
    p = tf.add_paragraph()
    p.space_before = Pt(space_before)
    p.space_after = Pt(3)
    p.level = level
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = FONT
    return p


def add_rounded_rect(slide, left, top, width, height, fill_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Emu(left), Emu(top), Emu(width), Emu(height))
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    return shape


def add_code_block(slide, left, top, width, height, code, size=9):
    rect = add_rounded_rect(slide, left, top, width, height, fill_color=LIGHT_BG)
    rect.text_frame.word_wrap = True
    p = rect.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = code
    r.font.size = Pt(size)
    r.font.name = "Courier New"
    r.font.color.rgb = DK1
    return rect


def set_slide_title(slide, title):
    """Set the title via the placeholder idx=15 that exists in layout Vide."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 15:
            ph.text = ""
            p = ph.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = title
            r.font.size = Pt(22)
            r.font.bold = True
            r.font.color.rgb = DK1
            r.font.name = FONT
            return
    add_textbox(slide, MARGIN_LEFT, 100000, 10800000, 700000,
                title, size=22, bold=True, color=DK1)


# ───────────────────────────────────────────────
# Build
# ───────────────────────────────────────────────

prs = Presentation(TEMPLATE)

while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get(
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

L_VIDE = prs.slide_layouts[0]
L_COVER = prs.slide_layouts[5]


def section_slide(number, title, subtitle=""):
    s = prs.slides.add_slide(L_VIDE)
    for ph in s.placeholders:
        if ph.placeholder_format.idx == 15:
            ph.text = ""
    add_textbox(s, MARGIN_LEFT, 1500000, 3000000, 2000000,
                number, size=120, bold=True, color=ACCENT1)
    add_textbox(s, 3800000, 1800000, 7800000, 1500000,
                title, size=44, bold=True, color=DK1)
    if subtitle:
        add_textbox(s, 3800000, 3500000, 7800000, 800000,
                    subtitle, size=18, color=GRAY)
    return s


def content_slide(title, bullets, code=None):
    s = prs.slides.add_slide(L_VIDE)
    set_slide_title(s, title)

    has_code = code is not None
    bw = 5600000 if has_code else 10800000

    if bullets:
        _, tf = add_textbox(s, MARGIN_LEFT, CONTENT_TOP, bw, 5500000, "", size=13)
        tf.paragraphs[0].clear()
        for b in bullets:
            if b.startswith("##"):
                add_para(tf, b[2:].strip(), size=14, bold=True, color=DK1, space_before=10)
            elif b.startswith(">>"):
                add_para(tf, "  " + b[2:].strip(), size=11, color=GRAY, level=1)
            else:
                add_para(tf, "• " + b, size=12, color=DK1)

    if has_code:
        add_code_block(s, 6600000, CONTENT_TOP, 5200000, 5500000, code, size=9)
    return s


def exercise_slide(number, title, instructions, code=None):
    """Slide dediee a un exercice avec badge orange."""
    s = prs.slides.add_slide(L_VIDE)
    set_slide_title(s, f"Exercice {number}")

    badge = add_rounded_rect(s, MARGIN_LEFT, CONTENT_TOP, 2200000, 380000,
                             fill_color=ACCENT2)
    badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = badge.text_frame.paragraphs[0].add_run()
    r.text = f"EXERCICE {number}"
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = FONT

    add_textbox(s, 3200000, CONTENT_TOP, 8400000, 380000,
                title, size=18, bold=True, color=DK1)

    has_code = code is not None
    iw = 5600000 if has_code else 10800000
    inst_top = CONTENT_TOP + 500000

    _, tf = add_textbox(s, MARGIN_LEFT, inst_top, iw, 5000000, "", size=12)
    tf.paragraphs[0].clear()
    for i, inst in enumerate(instructions):
        if inst.startswith(">>"):
            add_para(tf, "  " + inst[2:].strip(), size=11, color=GRAY, level=1)
        else:
            add_para(tf, f"{i+1}. {inst}" if not inst.startswith("-") else inst,
                     size=12, color=DK1, space_before=4)

    if has_code:
        add_code_block(s, 6600000, inst_top, 5200000, 5000000, code, size=9)
    return s


# ═══════════════════════════════════════════════
# COUVERTURE
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(L_COVER)
add_textbox(s, 400000, 1000000, 11400000, 1200000,
            "Formation dbt avec Snowflake",
            size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, 400000, 2400000, 11400000, 800000,
            "De Debutant a Expert",
            size=28, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, 400000, 3800000, 11400000, 500000,
            "Duree : 2 jours (16h)  |  Dataset : Snowflake TPCH_SF1",
            size=16, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, 400000, 5800000, 11400000, 500000,
            "KPC Group", size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# TABLE DES MATIERES
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(L_VIDE)
set_slide_title(s, "Programme de la formation")

modules = [
    ("01", "Introduction a dbt", "ELT vs ETL, stack data, dataset TPCH"),
    ("02", "Mise en place", "Installation, Snowflake, dbt init, GitHub"),
    ("03", "Construction du projet", "Sources, staging, intermediate, marts"),
    ("04", "Tests avances", "Severite, dbt_utils, dbt_expectations, unit tests"),
    ("05", "Modelisation avancee", "Seeds, snapshots, incremental, governance"),
    ("06", "Commandes et selecteurs", "CLI, selecteurs de graphe, tags, docs"),
    ("07", "Jinja et Macros", "Templates, macros, hooks, boucles"),
    ("08", "Deploiement dbt Cloud", "Service account, jobs, CI/CD"),
]

y = CONTENT_TOP + 100000
for i, (num, title, desc) in enumerate(modules):
    if i == 0:
        add_textbox(s, MARGIN_LEFT, y - 250000, 10800000, 200000,
                    "JOUR 1", size=11, bold=True, color=ACCENT2)
    if i == 4:
        add_textbox(s, MARGIN_LEFT, y - 250000, 10800000, 200000,
                    "JOUR 2", size=11, bold=True, color=ACCENT2)

    badge = add_rounded_rect(s, MARGIN_LEFT, y, 500000, 350000, fill_color=ACCENT1)
    badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = badge.text_frame.paragraphs[0].add_run()
    r.text = num
    r.font.size = Pt(15)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = FONT

    add_textbox(s, 1500000, y, 4200000, 350000, title, size=15, bold=True, color=DK1)
    add_textbox(s, 5900000, y + 30000, 5600000, 350000, desc, size=11, color=GRAY)
    y += 600000


# ═══════════════════════════════════════════════
# MODULE 1 : Introduction
# ═══════════════════════════════════════════════
section_slide("01", "Introduction a dbt", "1h00 — Theorie et decouverte")

content_slide("Qu'est-ce que dbt ?", [
    "dbt = data build tool : outil de transformation (le T de ELT)",
    "Les analystes ecrivent des SELECT, dbt gere le DDL/DML",
    "Approche ELT : Extract → Load → Transform (dans le warehouse)",
    "##dbt Core vs dbt Cloud",
    "dbt Core : open source, CLI, gratuit",
    "dbt Cloud : plateforme SaaS, IDE web, scheduling, CI/CD",
])

content_slide("Pourquoi dbt ?", [
    "##Les 5 benefices cles",
    "Inference des dependances : ref() construit le DAG automatiquement",
    "Tests integres : unique, not_null, relationships, accepted_values",
    "Documentation : site auto-genere depuis les fichiers YAML",
    "Version control : modeles SQL dans Git = code review + CI/CD",
    "Modularite : macros Jinja, packages reutilisables",
    "##Place dans la stack data",
    "Sources → Ingestion (Fivetran) → Snowflake Raw → dbt → Analytics → BI",
])

content_slide("Le dataset TPCH_SF1", [
    "##Schema de donnees",
    "REGION (5) → NATION (25) → CUSTOMER (150K)",
    "CUSTOMER → ORDERS (1.5M) → LINEITEM (6M)",
    "SUPPLIER (10K) → PARTSUPP (800K) → PART (200K)",
    "##Domaine metier",
    "Distributeur de pieces : clients, commandes, fournisseurs",
],
    code="-- Explorer les donnees\nSELECT 'CUSTOMER' AS tbl,\n       COUNT(*) AS cnt\nFROM SNOWFLAKE_SAMPLE_DATA\n     .TPCH_SF1.CUSTOMER\nUNION ALL\nSELECT 'ORDERS', COUNT(*)\nFROM SNOWFLAKE_SAMPLE_DATA\n     .TPCH_SF1.ORDERS\nUNION ALL\nSELECT 'LINEITEM', COUNT(*)\nFROM SNOWFLAKE_SAMPLE_DATA\n     .TPCH_SF1.LINEITEM\nORDER BY cnt DESC;")

content_slide("Architecture du projet en 3 couches", [
    "##Staging (vues)",
    "1 modele par table source",
    "Renommage, cast, nettoyage",
    "Convention : stg_tpch__<table>",
    "##Intermediate (ephemeral)",
    "Jointures et enrichissements",
    "Pas de table creee dans Snowflake",
    "Convention : int_<description>",
    "##Marts (tables)",
    "Modeles business-ready pour la BI",
    "Convention : fct_, dim_, agg_",
])


# ═══════════════════════════════════════════════
# MODULE 2 : Mise en place
# ═══════════════════════════════════════════════
section_slide("02", "Mise en place de\nl'environnement", "1h30 — Installation et configuration")

content_slide("Installation des outils", [
    "##Pre-requis",
    "Python 3.10+",
    "VS Code + extensions (dbt Power User, SQLFluff, YAML)",
    "Git + compte GitHub",
    "Acces Snowflake verifie",
],
    code="# Creer le projet\nmkdir tpch_analytics\ncd tpch_analytics\n\n# Environnement virtuel\npython3 -m venv venv\nsource venv/bin/activate\n\n# Installer dbt\npip install dbt-snowflake\n\n# Verifier\ndbt --version")

content_slide("Configuration Snowflake", [
    "##profiles.yml (~/.dbt/)",
    "Connexion SSO avec externalbrowser",
    "Schema dynamique par developpeur",
    "Ne JAMAIS commiter ce fichier !",
    "##Preparation Snowflake",
    "Creer la base ANALYTICS",
    "Creer les schemas STAGING, MARTS, SNAPSHOTS",
],
    code="# ~/.dbt/profiles.yml\ntpch_analytics:\n  target: dev\n  outputs:\n    dev:\n      type: snowflake\n      account: \"<ACCOUNT>\"\n      user: \"<USER>\"\n      authenticator: externalbrowser\n      role: ACCOUNTADMIN\n      database: ANALYTICS\n      warehouse: COMPUTE_WH\n      schema: DEV_{{ env_var('DBT_USER') }}\n      threads: 4\n\n-- Dans Snowflake :\nCREATE DATABASE IF NOT EXISTS ANALYTICS;\nCREATE SCHEMA ANALYTICS.STAGING;\nCREATE SCHEMA ANALYTICS.MARTS;")

content_slide("dbt_project.yml et structure", [
    "##Configuration centrale",
    "Materialisation par couche (view/ephemeral/table)",
    "Schemas personnalises (staging, marts)",
    "Tags pour categoriser les modeles",
],
    code="# dbt_project.yml\nname: 'tpch_analytics'\nversion: '1.0.0'\nprofile: 'tpch_analytics'\n\nmodels:\n  tpch_analytics:\n    staging:\n      +materialized: view\n      +schema: staging\n    intermediate:\n      +materialized: ephemeral\n    marts:\n      +materialized: table\n      +schema: marts\n      +tags: ['daily']")

exercise_slide("2A", "Verifier la connexion Snowflake", [
    "Ouvrir un terminal dans le dossier tpch_analytics/",
    "Activer l'environnement virtuel : source venv/bin/activate",
    "Lancer dbt debug",
    "Verifier que tous les checks sont verts",
    ">>Resultat attendu : Connection OK, All checks passed",
],
    code="$ cd tpch_analytics\n$ source venv/bin/activate\n$ dbt debug\n\n# Resultat attendu :\n# Connection:\n#   account: QRQOCXI-KD48953\n#   database: ANALYTICS\n#   schema: DEV_JULIEN\n#   warehouse: COMPUTE_WH\n#   All checks passed!")

exercise_slide("2B", "Initialiser le depot Git", [
    "Creer le fichier .gitignore",
    "Initialiser git dans le dossier projet",
    "Faire le commit initial",
    "Creer un repo sur GitHub",
    "Pousser le code",
],
    code="$ git init\n\n# .gitignore\ntarget/\ndbt_packages/\nlogs/\nvenv/\nprofiles.yml\n.env\n\n$ git add .\n$ git commit -m \"feat: initial dbt project\"\n$ git remote add origin <URL>\n$ git push -u origin main")


# ═══════════════════════════════════════════════
# MODULE 3 : Construction du projet
# ═══════════════════════════════════════════════
section_slide("03", "Construction du\nprojet dbt", "3h00 — Sources, staging, intermediate, marts")

content_slide("Declaration des sources", [
    "##models/sources.yml",
    "Declare les tables TPCH_SF1 comme sources",
    "Pointe vers SNOWFLAKE_SAMPLE_DATA",
    "##Principe cle",
    "Ne JAMAIS hardcoder les noms de table",
    "Utiliser {{ source('tpch', 'ORDERS') }}",
    "Remplacement de source en un seul endroit",
],
    code="# models/sources.yml\nversion: 2\nsources:\n  - name: tpch\n    database: SNOWFLAKE_SAMPLE_DATA\n    schema: TPCH_SF1\n    tables:\n      - name: ORDERS\n        columns:\n          - name: O_ORDERKEY\n            tests:\n              - unique\n              - not_null\n      - name: CUSTOMER\n      - name: LINEITEM\n      - name: PART\n      - name: SUPPLIER\n      - name: PARTSUPP\n      - name: NATION\n      - name: REGION")

content_slide("Modeles staging — Exemple", [
    "##Convention",
    "1 modele = 1 table source",
    "CTE source → CTE renamed → select *",
    "Renommage snake_case, noms metier",
    "Colonnes calculees si besoin",
    "##stg_tpch__line_items.sql",
    "Cle surrogate via dbt_utils",
    "Prix remise et prix final calcules",
],
    code="-- stg_tpch__orders.sql\nwith source as (\n    select *\n    from {{ source('tpch', 'ORDERS') }}\n),\nrenamed as (\n    select\n        o_orderkey    as order_key,\n        o_custkey     as customer_key,\n        o_orderstatus as order_status,\n        o_totalprice  as total_price,\n        o_orderdate   as order_date,\n        o_orderpriority\n                      as order_priority\n    from source\n)\nselect * from renamed")

content_slide("Tests et documentation staging", [
    "##4 tests generiques dbt",
    "unique : aucun doublon",
    "not_null : aucune valeur NULL",
    "accepted_values : valeurs dans une liste",
    "relationships : integrite referentielle",
    "##Fichier _stg__models.yml",
    "1 entry par modele staging",
    "Tests sur les cles et colonnes critiques",
],
    code="# _stg__models.yml\nversion: 2\nmodels:\n  - name: stg_tpch__orders\n    columns:\n      - name: order_key\n        tests:\n          - unique\n          - not_null\n      - name: customer_key\n        tests:\n          - not_null\n          - relationships:\n              to: ref('stg_tpch__customers')\n              field: customer_key\n      - name: order_status\n        tests:\n          - accepted_values:\n              values: ['F', 'O', 'P']")

exercise_slide("3A", "Creer les modeles staging manquants", [
    "Creer stg_tpch__parts.sql (part_key, part_name, brand...)",
    "Creer stg_tpch__suppliers.sql (supplier_key, supplier_name...)",
    "Creer stg_tpch__part_suppliers.sql (avec surrogate key)",
    "Creer stg_tpch__nations.sql (nation_key, nation_name...)",
    "Creer stg_tpch__regions.sql (region_key, region_name...)",
    "Executer : dbt run --select staging",
    "Tester : dbt test --select staging",
    ">>Indice : pour part_suppliers, utiliser generate_surrogate_key(['ps_partkey', 'ps_suppkey'])",
])

content_slide("Packages dbt", [
    "##packages.yml",
    "dbt_utils : macros utilitaires (surrogate_key, date_spine...)",
    "dbt_expectations : tests avances (Great Expectations style)",
    "##Installation",
    ">>dbt deps",
],
    code="# packages.yml\npackages:\n  - package: dbt-labs/dbt_utils\n    version: \">=1.3.0\"\n  - package: calogica/dbt_expectations\n    version: \">=0.10.0\"\n\n$ dbt deps\n# Installing dbt-labs/dbt_utils\n# Installing calogica/dbt_expectations\n# Packages installed!")

content_slide("Modeles intermediaires", [
    "##int_order_items.sql",
    "Jointure commandes + lignes de commande",
    "##int_order_items_summary.sql",
    "Agregation au niveau commande",
    ">>sum(quantity), sum(final_price), count(*)",
    "##int_part_suppliers_agg.sql",
    "Agregation pieces par fournisseur",
    "##Rappel : ephemeral",
    ">>Pas de table/vue creee dans Snowflake",
],
    code="-- int_order_items.sql\nwith orders as (\n    select * from {{ ref(\n      'stg_tpch__orders') }}\n),\nline_items as (\n    select * from {{ ref(\n      'stg_tpch__line_items') }}\n),\njoined as (\n    select\n        li.line_item_key,\n        li.order_key,\n        o.customer_key,\n        o.order_date,\n        o.order_status,\n        li.quantity,\n        li.final_price,\n        li.ship_date,\n        li.return_flag\n    from line_items li\n    inner join orders o\n      on li.order_key\n       = o.order_key\n)\nselect * from joined")

content_slide("Modeles marts", [
    "##fct_orders.sql (table de faits)",
    "Enrichie avec nom client, segment, days_to_fulfill",
    "##dim_customers / dim_parts / dim_suppliers",
    "Dimensions enrichies avec geographie et metriques",
    "##agg_customer_lifetime_value.sql",
    "total_orders, lifetime_value, tenure_days",
    "##agg_supplier_performance.sql",
    "revenue, return_rate, avg_delivery_days",
],
    code="-- fct_orders.sql\nwith order_summary as (\n    select * from {{ ref(\n      'int_order_items_summary') }}\n),\ncustomers as (\n    select * from {{ ref(\n      'stg_tpch__customers') }}\n),\nfinal as (\n    select\n        os.order_key,\n        c.customer_name,\n        c.market_segment,\n        os.order_date,\n        os.net_amount,\n        os.line_item_count,\n        datediff('day',\n          os.order_date,\n          os.last_ship_date)\n          as days_to_fulfill\n    from order_summary os\n    join customers c\n      on os.customer_key\n       = c.customer_key\n)\nselect * from final")

content_slide("Materialisations", [
    "##view (staging)",
    "Vue SQL, recalculee a chaque requete",
    "##table (marts)",
    "Table physique, recalculee a chaque dbt run",
    "##ephemeral (intermediate)",
    "Pas de table/vue, SQL inline via CTE",
    "##incremental (Module 5)",
    "Table avec ajout/mise a jour incrementale",
])

exercise_slide("3C", "Build complet et exploration", [
    "Lancer : dbt build",
    "Verifier que tous les modeles et tests passent",
    "Dans Snowflake, requeter les marts :",
    ">>Top 10 clients par lifetime_value",
    ">>Fournisseurs avec les plus hauts taux de retour",
    ">>Repartition des commandes par segment de marche",
],
    code="$ dbt build\n\n-- Puis dans Snowflake :\nSELECT customer_name,\n       market_segment,\n       lifetime_value\nFROM ANALYTICS.MARTS\n  .AGG_CUSTOMER_LIFETIME_VALUE\nORDER BY lifetime_value DESC\nLIMIT 10;\n\nSELECT supplier_name,\n       total_revenue,\n       return_rate\nFROM ANALYTICS.MARTS\n  .AGG_SUPPLIER_PERFORMANCE\nWHERE orders_fulfilled > 100\nORDER BY return_rate DESC\nLIMIT 10;")

exercise_slide("3D", "Changer la materialisation", [
    "Modifier int_order_items.sql : ajouter {{ config(materialized='view') }}",
    "Lancer : dbt run --select int_order_items",
    "Observer dans Snowflake : la vue INT_ORDER_ITEMS apparait",
    "Remettre en ephemeral (supprimer le config)",
    "Relancer : dbt run --select int_order_items+",
])


# ═══════════════════════════════════════════════
# MODULE 4 : Tests avances
# ═══════════════════════════════════════════════
section_slide("04", "Tests avances", "1h30 — Qualite des donnees")

content_slide("Severite et seuils des tests", [
    "##Niveaux de severite",
    "error : bloque le dbt build",
    "warn : avertissement, ne bloque pas",
    "##Seuils configurables",
    "warn_if: \">10\" → warn si >10 erreurs",
    "error_if: \">100\" → error si >100 erreurs",
    "##Configuration globale",
    ">>dbt_project.yml : tests.+severity: warn",
],
    code="# Dans _stg__models.yml\n- name: order_key\n  tests:\n    - unique:\n        severity: error\n    - not_null:\n        severity: warn\n        warn_if: \">10\"\n        error_if: \">100\"\n\n# Configuration globale\n# dbt_project.yml\ntests:\n  tpch_analytics:\n    +severity: warn")

content_slide("Tests dbt_utils", [
    "##unique_combination_of_columns",
    "Combinaison unique de plusieurs colonnes",
    "##accepted_range",
    "Valeur dans un intervalle [min, max]",
    "##expression_is_true",
    "Expression SQL personnalisee",
],
    code="# _stg__models.yml\nmodels:\n  - name: stg_tpch__line_items\n    tests:\n      - dbt_utils\n        .unique_combination_of_columns:\n          combination_of_columns:\n            - order_key\n            - line_number\n    columns:\n      - name: discount_percentage\n        tests:\n          - dbt_utils.accepted_range:\n              min_value: 0\n              max_value: 1\n              inclusive: true\n      - name: ship_date\n        tests:\n          - dbt_utils.expression_is_true:\n              expression: \">= '1992-01-01'\"")

content_slide("Tests dbt_expectations", [
    "##expect_column_values_to_be_between",
    "Avec condition optionnelle (row_condition)",
    "##expect_column_values_to_be_of_type",
    "Verification du type de colonne",
    "##expect_table_row_count_to_be_between",
    "Nombre de lignes dans un intervalle",
],
    code="# _marts__models.yml\nmodels:\n  - name: fct_orders\n    columns:\n      - name: net_amount\n        tests:\n          - dbt_expectations\n            .expect_column_values_to_be_between:\n              min_value: 0\n              max_value: 1000000\n              row_condition: \"order_status != 'F'\"\n      - name: order_date\n        tests:\n          - dbt_expectations\n            .expect_column_values_to_be_of_type:\n              column_type: date\n    tests:\n      - dbt_expectations\n        .expect_table_row_count_to_be_between:\n          min_value: 1000000\n          max_value: 2000000")

content_slide("Tests personnalises", [
    "##Tests singuliers (tests/singular/)",
    "Fichier SQL qui retourne 0 ligne si OK",
    "##Tests generiques (tests/generic/)",
    "Macro Jinja reutilisable sur n'importe quel modele",
],
    code="-- tests/singular/\n--   assert_positive_order_totals.sql\nselect order_key, net_amount\nfrom {{ ref('fct_orders') }}\nwhere net_amount < 0\n\n-- tests/generic/\n--   test_accepted_range.sql\n{% test accepted_range(\n    model, column_name,\n    min_val, max_val) %}\nselect {{ column_name }}\nfrom {{ model }}\nwhere {{ column_name }} < {{ min_val }}\n   or {{ column_name }} > {{ max_val }}\n{% endtest %}\n\n# Utilisation dans le YAML :\n- name: quantity\n  tests:\n    - accepted_range:\n        min_val: 1\n        max_val: 50")

content_slide("Unit tests (dbt 1.8+)", [
    "##Principe",
    "Tester la logique avec des donnees fictives",
    "given : inputs mock (ref simulees)",
    "expect : resultat attendu",
    "##Avantages",
    "Pas besoin de donnees reelles",
    "Test de la logique SQL pure",
    "Rapide a executer",
],
    code="# _marts__models.yml\nunit_tests:\n  - name: test_clv_calculation\n    model: agg_customer_lifetime_value\n    given:\n      - input: ref('fct_orders')\n        rows:\n          - {order_key: 100,\n             customer_key: 1,\n             net_amount: 500.00,\n             order_date: \"1995-01-01\", ...}\n          - {order_key: 101,\n             customer_key: 1,\n             net_amount: 300.00,\n             order_date: \"1995-06-01\", ...}\n      - input: ref('dim_customers')\n        rows:\n          - {customer_key: 1, ...}\n    expect:\n      rows:\n        - {customer_key: 1,\n           total_orders: 2,\n           lifetime_value: 800.00}\n\n$ dbt test --select test_type:unit")

exercise_slide("4A", "Ecrire un test singulier", [
    "Creer tests/singular/assert_positive_lifetime_value.sql",
    "Le test doit verifier qu'aucun client n'a une lifetime_value negative",
    "Lancer : dbt test --select assert_positive_lifetime_value",
    ">>Le test doit passer (0 lignes retournees)",
])

exercise_slide("4B", "Ecrire un unit test", [
    "Ajouter un unit test pour agg_supplier_performance",
    "Verifier le calcul du return_rate",
    "Mock : 2 lignes dont 1 avec return_flag = 'R'",
    "Attendu : return_rate = 0.5",
    "Lancer : dbt test --select test_type:unit",
])


# ═══════════════════════════════════════════════
# MODULE 5 : Modelisation avancee
# ═══════════════════════════════════════════════
section_slide("05", "Modelisation\navancee", "2h00 — Seeds, snapshots, incremental, governance")

content_slide("Seeds et fonction doc()", [
    "##Seeds = fichiers CSV → tables",
    "Chargement via dbt seed",
    "Tables de reference statiques",
    "##Fonction {{ doc() }}",
    "Documentation Markdown dans fichiers .md",
    "Reference dans le YAML : {{ doc('order_status') }}",
],
    code="# seeds/order_priority_mapping.csv\norder_priority,priority_label,rank\n1-URGENT,Urgent,1\n2-HIGH,High,2\n3-MEDIUM,Medium,3\n4-NOT SPECIFIED,Not Specified,4\n5-LOW,Low,5\n\n$ dbt seed\n\n# _stg__docs.md\n{% docs order_status %}\n- **F** = Fulfilled\n- **O** = Open\n- **P** = Partial\n{% enddocs %}\n\n# YAML : {{ doc('order_status') }}")

exercise_slide("5A", "Creer un seed region_mapping", [
    "Creer seeds/region_mapping.csv avec les colonnes :",
    ">>region_name, business_region",
    "Mapper : AFRICA/EUROPE/MIDDLE EAST → EMEA",
    "Mapper : AMERICA → AMERICAS, ASIA → APAC",
    "Charger : dbt seed",
    "Enrichir dim_customers.sql avec business_region",
    "Lancer : dbt run --select dim_customers",
])

content_slide("Snapshots — SCD Type 2", [
    "##Principe",
    "Capturer les changements au fil du temps",
    "Slowly Changing Dimension Type 2",
    "##Strategies",
    "check : compare les colonnes specifiees",
    "timestamp : colonne de date de modification",
    "##Colonnes SCD2 ajoutees automatiquement",
    "dbt_valid_from, dbt_valid_to, dbt_scd_id",
],
    code="-- snapshots/snap_customers.sql\n{% snapshot snap_customers %}\n{{\n    config(\n        target_database='ANALYTICS',\n        target_schema='SNAPSHOTS',\n        unique_key='customer_key',\n        strategy='check',\n        check_cols=[\n          'account_balance',\n          'market_segment',\n          'address', 'phone'\n        ],\n    )\n}}\nselect *\nfrom {{ ref('stg_tpch__customers') }}\n{% endsnapshot %}\n\n$ dbt snapshot")

exercise_slide("5B", "Executer le snapshot 2 fois", [
    "Lancer : dbt snapshot (1ere fois)",
    "Inspecter ANALYTICS.SNAPSHOTS.SNAP_CUSTOMERS dans Snowflake",
    "Observer les colonnes dbt_valid_from, dbt_valid_to",
    "Lancer : dbt snapshot (2eme fois)",
    "Observer : pas de nouvelles lignes (donnees inchangees)",
    "Discussion : que se passerait-il si market_segment changeait ?",
],
    code="-- Inspecter le resultat\nSELECT\n    customer_key,\n    customer_name,\n    market_segment,\n    dbt_valid_from,\n    dbt_valid_to,\n    dbt_updated_at\nFROM ANALYTICS.SNAPSHOTS\n     .SNAP_CUSTOMERS\nWHERE customer_key = 1;")

content_slide("Modeles incrementaux", [
    "##Principe",
    "Ne traite que les nouvelles donnees",
    "##Concepts cles",
    "is_incremental() : table existe deja ?",
    "{{ this }} : table cible existante",
    "--full-refresh : reconstruction complete",
    "##Strategies Snowflake",
    "merge (UPSERT), delete+insert, append",
],
    code="-- fct_order_items.sql\n{{\n    config(\n        materialized='incremental',\n        unique_key='line_item_key',\n        incremental_strategy='merge'\n    )\n}}\nwith order_items as (\n    select *\n    from {{ ref('int_order_items') }}\n)\nselect *\nfrom order_items\n{% if is_incremental() %}\n  where ship_date > (\n    select max(ship_date)\n    from {{ this }}\n  )\n{% endif %}\n\n-- Premier run :\n$ dbt run --select fct_order_items\n    --full-refresh\n-- Runs suivants :\n$ dbt run --select fct_order_items")

content_slide("Model Governance", [
    "##Groupes et acces (dbt 1.5+)",
    "Groupes avec owners (equipe, email)",
    "Niveaux : public / protected / private",
    "##Contrats",
    "Garantir le schema de sortie (data_type obligatoire)",
    "dbt run echoue si schema different",
],
    code="# dbt_project.yml\ngroups:\n  - name: finance\n    owner:\n      name: \"Equipe Finance\"\n      email: \"finance@co.com\"\n\n# _marts__models.yml\n- name: fct_orders\n  group: finance\n  access: public\n  config:\n    contract:\n      enforced: true\n  columns:\n    - name: order_key\n      data_type: number(38,0)\n    - name: net_amount\n      data_type: number(12,2)")


# ═══════════════════════════════════════════════
# MODULE 6 : Commandes et selecteurs
# ═══════════════════════════════════════════════
section_slide("06", "Commandes et\nselecteurs", "1h00 — Maitriser la CLI dbt")

content_slide("Commandes essentielles", [
    "##Cycle de developpement",
    "dbt debug / dbt deps / dbt seed",
    "dbt run / dbt test / dbt build",
    "dbt compile / dbt snapshot / dbt clean",
    "##Documentation",
    "dbt docs generate → dbt docs serve",
    "##Autres",
    "dbt source freshness",
    "dbt ls (lister les ressources)",
    "dbt run-operation <macro>",
])

content_slide("Selecteurs de graphe", [
    "##Operateurs",
    "model+ → modele et TOUS ses descendants",
    "+model → TOUS ses ancetres et le modele",
    "model+1 → enfants directs (1 niveau)",
    "staging → tout le dossier",
    "tag:daily → par tag",
    "--exclude model → exclure",
],
    code="# Exemples\ndbt run --select stg_tpch__orders+\ndbt run --select +fct_orders\ndbt run --select staging\ndbt run --select stg_tpch__orders+1\n\ndbt run --select marts \\\n  --exclude agg_supplier_performance\n\ndbt run --select tag:daily\ndbt test --select tag:critical\n\n# Lister les ressources\ndbt ls --select \\\n  +agg_customer_lifetime_value\n\n# Compiler pour voir le SQL\ndbt compile --select fct_orders")

exercise_slide("6A", "Challenges selecteurs", [
    "Lister tous les ancetres de agg_customer_lifetime_value",
    ">>dbt ls --select +agg_customer_lifetime_value",
    "Executer intermediate et leurs marts directs",
    ">>dbt run --select intermediate+1",
    "Tester uniquement les marts",
    ">>dbt test --select marts",
    "Compiler fct_orders et inspecter le SQL genere",
    ">>dbt compile --select fct_orders",
    ">>Regarder target/compiled/.../fct_orders.sql",
])

exercise_slide("6B", "Explorer la documentation", [
    "Generer la doc : dbt docs generate",
    "Lancer le serveur : dbt docs serve",
    "Explorer le graphe de lignage (DAG)",
    "Verifier que tous les modeles ont des descriptions",
    "Inspecter la page de fct_orders",
])


# ═══════════════════════════════════════════════
# MODULE 7 : Jinja et Macros
# ═══════════════════════════════════════════════
section_slide("07", "Jinja et Macros", "2h00 — Templates, macros, hooks")

content_slide("Fondamentaux Jinja", [
    "##3 types de blocs",
    "{# commentaire #} — ignore dans le SQL",
    "{% statement %} — logique (set, for, if)",
    "{{ expression }} — affichage de valeur",
    "##Variables et boucles",
    "{% set var = 'hello' %}",
    "{% for item in list %} ... {% endfor %}",
    "##Astuce",
    ">>dbt compile pour voir le SQL genere",
],
    code="{% set price_columns = [\n  'extended_price',\n  'discounted_price',\n  'final_price'\n] %}\n\nselect\n    order_key,\n    {% for col in price_columns %}\n        sum({{ col }})\n          as total_{{ col }}\n        {% if not loop.last %},\n        {% endif %}\n    {% endfor %}\nfrom {{ ref('int_order_items') }}\ngroup by 1\n\n-- SQL genere :\n-- sum(extended_price)\n--   as total_extended_price,\n-- sum(discounted_price)\n--   as total_discounted_price,\n-- sum(final_price)\n--   as total_final_price")

content_slide("Construire une macro", [
    "##Macros = fonctions reutilisables",
    "Fichier dans macros/",
    "Appelees dans les modeles SQL",
    "##generate_schema_name",
    "Controle le schema de destination",
    "Dev : DEV_JULIEN_staging",
    "Prod : staging (directement)",
],
    code="-- macros/discount_amount.sql\n{% macro discount_amount(\n    price_col, discount_col) %}\n  ({{ price_col }} * {{ discount_col }})\n{% endmacro %}\n\n-- macros/generate_schema_name.sql\n{% macro generate_schema_name(\n    custom_schema_name, node) -%}\n  {%- set default = target.schema -%}\n  {%- if custom_schema_name is none -%}\n    {{ default }}\n  {%- elif target.name == 'prod' -%}\n    {{ custom_schema_name | trim }}\n  {%- else -%}\n    {{ default }}_{{ custom_schema_name }}\n  {%- endif -%}\n{%- endmacro %}")

content_slide("Hooks et operations", [
    "##Hooks = SQL avant/apres un modele",
    "post-hook : GRANT, log, notification",
    "pre-hook : nettoyage, preparation",
    "##Operations = macros standalone",
    ">>dbt run-operation log_run_start",
    "##Variable target",
    "target.name, target.schema, target.database",
    "Limiter les donnees en dev",
],
    code="-- macros/grant_select.sql\n{% macro grant_select(\n    role='ANALYST') %}\n  {% set sql %}\n    GRANT SELECT ON {{ this }}\n      TO ROLE {{ role }};\n  {% endset %}\n  {% do run_query(sql) %}\n{% endmacro %}\n\n# dbt_project.yml\nmodels:\n  tpch_analytics:\n    marts:\n      +post-hook:\n        - \"{{ grant_select('ANALYST') }}\"\n\n-- Limiter en dev :\nselect *\nfrom {{ ref('int_order_items') }}\n{% if target.name == 'dev' %}\n  where order_date >=\n    dateadd('year', -1,\n            current_date())\n{% endif %}")

exercise_slide("7A", "Creer la macro star_except", [
    "Creer macros/star_except.sql",
    "La macro selectionne toutes les colonnes sauf celles specifiees",
    "Utiliser adapter.get_columns_in_relation()",
    "Tester dans un modele :",
    ">>select {{ star_except(ref('stg_tpch__customers'), ['comment', 'address']) }}",
    ">>from {{ ref('stg_tpch__customers') }}",
    "Compiler pour verifier : dbt compile --select <modele>",
])

exercise_slide("7B", "Modele pivot avec boucle for", [
    "Creer analyses/monthly_orders_by_priority.sql",
    "Utiliser une boucle for pour generer une colonne par priorite",
    "Utiliser dbt_utils.get_column_values() pour les valeurs",
    "Resultat : un tableau pivot (lignes=mois, colonnes=priorites)",
    "Compiler pour verifier le SQL genere",
],
    code="{% set priorities = dbt_utils\n  .get_column_values(\n    table=ref('stg_tpch__orders'),\n    column='order_priority'\n  ) %}\n\nselect\n  date_trunc('month', order_date)\n    as order_month,\n  {% for p in priorities %}\n    count(case\n      when order_priority = '{{ p }}'\n      then 1 end)\n      as {{ p | replace('-','_') | lower }}\n    {% if not loop.last %},{% endif %}\n  {% endfor %}\nfrom {{ ref('stg_tpch__orders') }}\ngroup by 1\norder by 1")


# ═══════════════════════════════════════════════
# MODULE 8 : Deploiement
# ═══════════════════════════════════════════════
section_slide("08", "Deploiement et\ndbt Cloud", "2h00 — Production et CI/CD")

content_slide("Profil de production", [
    "##Compte de service Snowflake",
    "Utilisateur STREAMLIT avec cle RSA",
    "Authentification key-pair (pas de mot de passe)",
    "Role SYSADMIN (pas ACCOUNTADMIN)",
    "##Differences dev vs prod",
    "Dev : 4 threads, SSO, schema DEV_<user>",
    "Prod : 8 threads, key-pair, schema PROD",
],
    code="# profiles.yml (prod)\ntpch_analytics:\n  target: prod\n  outputs:\n    prod:\n      type: snowflake\n      account: QRQOCXI-KD48953\n      user: STREAMLIT\n      private_key_path: /path/key.p8\n      role: SYSADMIN\n      database: ANALYTICS\n      warehouse: COMPUTE_WH\n      schema: PROD\n      threads: 8")

content_slide("Setup dbt Cloud", [
    "##Etapes",
    "1. Creer un compte sur cloud.getdbt.com",
    "2. Connecter Snowflake (cle RSA du compte STREAMLIT)",
    "3. Connecter le repo GitHub",
    "4. Configurer : dbt version, target name, schema",
    "##L'IDE dbt Cloud",
    "Editeur de code avec coloration syntaxique",
    "Integration Git (branches, commits, PR)",
    "Graphe de lignage interactif",
])

content_slide("Configuration des jobs", [
    "##Job 1 : Run quotidien (06:00 UTC)",
    ">>dbt seed → dbt run → dbt test",
    "##Job 2 : Full refresh hebdo (dim 02:00)",
    ">>dbt run --full-refresh → dbt test",
    ">>dbt source freshness",
    "##Job 3 : CI/CD sur Pull Request",
    ">>dbt build --select state:modified+",
    "Slim CI : ne construit que les modeles modifies",
    "##Source freshness",
    "Verifie que les donnees sont a jour",
    ">>loaded_at_field + warn_after / error_after",
])

exercise_slide("8A", "Workflow complet (exercice final)", [
    "Creer une branche : git checkout -b feature/add-monthly-revenue",
    "Creer models/marts/agg_monthly_revenue.sql",
    "Ajouter tests et documentation dans _marts__models.yml",
    "Tester : dbt build --select agg_monthly_revenue+",
    "Commit et push vers GitHub",
    "Ouvrir une Pull Request",
    "Observer le job CI dans dbt Cloud",
    "Merger et observer le deploy production",
],
    code="-- agg_monthly_revenue.sql\nwith orders as (\n    select *\n    from {{ ref('fct_orders') }}\n),\nmonthly as (\n    select\n        date_trunc('month',\n          order_date) as order_month,\n        market_segment,\n        count(distinct order_key)\n          as total_orders,\n        sum(net_amount)\n          as total_revenue\n    from orders\n    group by 1, 2\n)\nselect * from monthly\n\n$ git add .\n$ git commit -m \"feat: monthly revenue\"\n$ git push -u origin \\\n    feature/add-monthly-revenue")


# ═══════════════════════════════════════════════
# AIDE-MEMOIRE
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(L_VIDE)
set_slide_title(s, "Aide-memoire des commandes dbt")

add_code_block(s, MARGIN_LEFT, CONTENT_TOP, 5400000, 5400000,
    "# Cycle de developpement\n"
    "dbt debug          # Connexion\n"
    "dbt deps           # Packages\n"
    "dbt seed           # CSV -> tables\n"
    "dbt run            # Modeles\n"
    "dbt test           # Tests\n"
    "dbt build          # Run+test (DAG)\n"
    "dbt compile        # SQL sans exec\n"
    "dbt snapshot       # Snapshots\n"
    "dbt clean          # Nettoyage\n\n"
    "# Documentation\n"
    "dbt docs generate\n"
    "dbt docs serve\n\n"
    "# Operations\n"
    "dbt run-operation <macro>\n"
    "dbt source freshness", size=10)

add_code_block(s, 6300000, CONTENT_TOP, 5400000, 5400000,
    "# Selecteurs\n"
    "dbt run --select <model>\n"
    "dbt run --select <model>+\n"
    "dbt run --select +<model>\n"
    "dbt run --select <folder>\n"
    "dbt run --select tag:<tag>\n"
    "dbt run --exclude <model>\n"
    "dbt run --full-refresh\n\n"
    "# Lister\n"
    "dbt ls --select <selector>\n\n"
    "# Tests par type\n"
    "dbt test --select test_type:unit\n"
    "dbt test --select test_type:singular\n\n"
    "# Materialisations\n"
    "view | table | ephemeral | incremental", size=10)


# ═══════════════════════════════════════════════
# MERCI
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(L_COVER)
add_textbox(s, 400000, 2000000, 11400000, 1500000,
            "Merci !", size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, 400000, 3800000, 11400000, 500000,
            "www.kpcgroup.fr", size=18, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
prs.save(OUTPUT)
print(f"OK: {OUTPUT}")
print(f"Slides: {len(prs.slides)}")
